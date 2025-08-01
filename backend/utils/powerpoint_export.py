from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image as PILImage
import os
import io
import base64
import requests
import tempfile
from typing import List, Dict, Any, Optional

# Global scale factor for canvas to PowerPoint conversion
CANVAS_TO_PPT_SCALE = 1.0

def _calculate_powerpoint_dimensions(canvas_width: int, canvas_height: int):
    """Calculate standard PowerPoint dimensions based on canvas aspect ratio"""
    # Standard PowerPoint formats:
    # - Standard (4:3): 10" x 7.5"  
    # - Widescreen (16:9): 13.33" x 7.5"
    
    aspect_ratio = canvas_width / canvas_height
    print(f"📏 Canvas aspect ratio: {aspect_ratio:.2f}")
    
    # Choose appropriate PowerPoint format based on aspect ratio
    if aspect_ratio > 1.6:  # Closer to 16:9 (1.78)
        slide_width_inches = 13.33
        slide_height_inches = 7.5
        print(f"📏 Using Widescreen format (16:9): {slide_width_inches}\" x {slide_height_inches}\"")
    else:  # Closer to 4:3 (1.33) or other ratios
        slide_width_inches = 10.0
        slide_height_inches = 7.5  
        print(f"📏 Using Standard format (4:3): {slide_width_inches}\" x {slide_height_inches}\"")
    
    # Calculate scale factor for element positioning
    global CANVAS_TO_PPT_SCALE
    CANVAS_TO_PPT_SCALE = min(slide_width_inches * 96 / canvas_width, slide_height_inches * 96 / canvas_height)
    print(f"📏 Canvas to PPT scale factor: {CANVAS_TO_PPT_SCALE:.3f}")
    
    return slide_width_inches, slide_height_inches

def create_unit_powerpoint(unit_data: Dict[str, Any], output_path: str, template_config: Optional[Dict[str, Any]] = None) -> str:
    """
    Create a PowerPoint presentation from a single naval unit
    
    Args:
        unit_data: Unit data including layout config
        output_path: Path where to save the PowerPoint file
        template_config: Optional template configuration for presentation format
    
    Returns:
        Path to the created PowerPoint file
    """
    
    try:
        print(f"Creating PowerPoint for unit: {unit_data.get('name', 'Unknown')}")
        
        # Create new presentation
        prs = Presentation()
        
        # Apply template configuration if provided
        print(f"🔍 Template config received: {template_config}")
        
        if template_config:
            canvas_width = template_config.get('canvasWidth', 1123)
            canvas_height = template_config.get('canvasHeight', 794)
            
            print(f"📏 Canvas dimensions from template: {canvas_width} x {canvas_height}")
            
            # Use standard PowerPoint dimensions
            slide_width_inches, slide_height_inches = _calculate_powerpoint_dimensions(canvas_width, canvas_height)
            
            prs.slide_width = Inches(slide_width_inches)
            prs.slide_height = Inches(slide_height_inches)
            
            print(f"✅ Applied PowerPoint standard dimensions: {slide_width_inches}\" x {slide_height_inches}\"")
        else:
            print("⚠️ No template config provided, using default")
            # Try to get canvas config from unit data
            unit_layout = unit_data.get('layout_config', {})
            if unit_layout:
                canvas_width = unit_layout.get('canvasWidth', 1123)
                canvas_height = unit_layout.get('canvasHeight', 794)
                
                print(f"📏 Using unit layout dimensions: {canvas_width} x {canvas_height}")
                
                # Apply same standard PowerPoint logic
                slide_width_inches, slide_height_inches = _calculate_powerpoint_dimensions(canvas_width, canvas_height)
                
                prs.slide_width = Inches(slide_width_inches)
                prs.slide_height = Inches(slide_height_inches)
                
                print(f"✅ Applied unit layout dimensions: {slide_width_inches:.2f}\" x {slide_height_inches:.2f}\"")
            else:
                # Default to 16:9 widescreen
                prs.slide_width = Inches(13.33)
                prs.slide_height = Inches(7.5)
                print(f"⚠️ Using default widescreen dimensions: 13.33\" x 7.5\"")
        
        # Create the unit slide
        slide = _create_unit_slide(prs, unit_data, {})
        
        # Save presentation
        print(f"Saving presentation to: {output_path}")
        prs.save(output_path)
        print(f"PowerPoint saved successfully")
        
        return output_path
        
    except Exception as e:
        print(f"Error in create_unit_powerpoint: {e}")
        import traceback
        traceback.print_exc()
        raise

def create_unit_powerpoint_to_buffer(unit_data: Dict[str, Any], output_buffer, template_config: Optional[Dict[str, Any]] = None) -> None:
    """
    Create a PowerPoint presentation from a single naval unit and save to buffer
    
    Args:
        unit_data: Unit data including layout config
        output_buffer: BytesIO buffer to save the PowerPoint
        template_config: Optional template configuration for presentation format
    """
    
    try:
        print(f"Creating PowerPoint for unit: {unit_data.get('name', 'Unknown')}")
        
        # Create new presentation
        prs = Presentation()
        
        # Apply template configuration if provided
        print(f"🔍 Template config received: {template_config}")
        
        if template_config:
            canvas_width = template_config.get('canvasWidth', 1123)
            canvas_height = template_config.get('canvasHeight', 794)
            
            print(f"📏 Canvas dimensions from template: {canvas_width} x {canvas_height}")
            
            # Use standard PowerPoint dimensions
            slide_width_inches, slide_height_inches = _calculate_powerpoint_dimensions(canvas_width, canvas_height)
            
            prs.slide_width = Inches(slide_width_inches)
            prs.slide_height = Inches(slide_height_inches)
            
            print(f"✅ Applied PowerPoint standard dimensions: {slide_width_inches}\" x {slide_height_inches}\"")
        else:
            print("⚠️ No template config provided, using default")
            # Try to get canvas config from unit data
            unit_layout = unit_data.get('layout_config', {})
            if unit_layout:
                canvas_width = unit_layout.get('canvasWidth', 1123)
                canvas_height = unit_layout.get('canvasHeight', 794)
                
                print(f"📏 Using unit layout dimensions: {canvas_width} x {canvas_height}")
                
                # Apply same standard PowerPoint logic
                slide_width_inches, slide_height_inches = _calculate_powerpoint_dimensions(canvas_width, canvas_height)
                
                prs.slide_width = Inches(slide_width_inches)
                prs.slide_height = Inches(slide_height_inches)
                
                print(f"✅ Applied unit layout dimensions: {slide_width_inches:.2f}\" x {slide_height_inches:.2f}\"")
            else:
                # Default to 16:9 widescreen
                prs.slide_width = Inches(13.33)
                prs.slide_height = Inches(7.5)
                print(f"⚠️ Using default widescreen dimensions: 13.33\" x 7.5\"")
        
        # Create the unit slide
        slide = _create_unit_slide(prs, unit_data, {})
        
        # Save presentation to buffer
        print(f"Saving presentation to buffer")
        prs.save(output_buffer)
        print(f"PowerPoint saved successfully to buffer")
        
    except Exception as e:
        print(f"Error in create_unit_powerpoint_to_buffer: {e}")
        import traceback
        traceback.print_exc()
        raise

def create_group_powerpoint(group_data: Dict[str, Any], output_path: str) -> str:
    """
    Create a PowerPoint presentation from a group's naval units
    
    Args:
        group_data: Group data including naval units and presentation config
        output_path: Path where to save the PowerPoint file
    
    Returns:
        Path to the created PowerPoint file
    """
    
    try:
        print(f"Creating PowerPoint for group: {group_data.get('name', 'Unknown')}")
        
        # Create new presentation
        prs = Presentation()
        
        # Set slide dimensions to 16:9 widescreen
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        print(f"Presentation created with widescreen dimensions")
        
        # Create title slide
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        title_slide = prs.slides.add_slide(title_slide_layout)
        
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        
        group_name = group_data.get('name', 'Gruppo Unità Navali')
        naval_units = group_data.get('naval_units', [])
        
        title.text = group_name
        subtitle.text = f"Presentazione di {len(naval_units)} unità navali"
        
        print(f"Title slide created for {len(naval_units)} units")
        
        # Get presentation config
        presentation_config = group_data.get('presentation_config', {})
        mode = presentation_config.get('mode', 'single')
        
        print(f"Presentation mode: {mode}")
        
        if mode == 'single':
            # Create one slide per unit
            for i, unit in enumerate(naval_units):
                print(f"Creating slide {i+1}/{len(naval_units)} for unit: {unit.get('name', 'Unknown')}")
                try:
                    slide = _create_unit_slide(prs, unit, group_data)
                except Exception as unit_error:
                    print(f"Failed to create slide for unit {unit.get('name', 'Unknown')}: {unit_error}")
                    # Continue with other units
                    continue
        else:
            # Create grid slides
            grid_rows = presentation_config.get('grid_rows', 3)
            grid_cols = presentation_config.get('grid_cols', 3)
            units_per_slide = grid_rows * grid_cols
            
            print(f"Creating grid slides: {grid_rows}x{grid_cols} ({units_per_slide} units per slide)")
            
            # Split units into pages
            for i in range(0, len(naval_units), units_per_slide):
                page_units = naval_units[i:i + units_per_slide]
                page_num = i // units_per_slide + 1
                print(f"Creating grid slide {page_num} with {len(page_units)} units")
                try:
                    slide = _create_grid_slide(prs, page_units, group_data, grid_rows, grid_cols, page_num)
                except Exception as grid_error:
                    print(f"Failed to create grid slide {page_num}: {grid_error}")
                    continue
        
        # Save presentation
        print(f"Saving presentation to: {output_path}")
        prs.save(output_path)
        print(f"PowerPoint saved successfully")
        
        return output_path
        
    except Exception as e:
        print(f"Error in create_group_powerpoint: {e}")
        import traceback
        traceback.print_exc()
        raise

def _create_unit_slide(prs: Presentation, unit: Dict[str, Any], group_data: Dict[str, Any]) -> Any:
    """Create a single slide for one naval unit"""
    
    try:
        print(f"Creating slide for unit: {unit.get('name', 'Unknown')}")
        
        # Use blank slide layout
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Get unit layout config with safe parsing
        layout_config = unit.get('layout_config', {})
        
        # Handle case where layout_config might be a JSON string
        if isinstance(layout_config, str):
            try:
                import json
                layout_config = json.loads(layout_config)
            except (json.JSONDecodeError, TypeError):
                print(f"Failed to parse layout_config JSON for unit {unit.get('name')}, using default")
                layout_config = {}
        
        elements = layout_config.get('elements', [])
        canvas_background = layout_config.get('canvasBackground', '#ffffff')
        
        print(f"Unit has {len(elements)} elements, background: {canvas_background}")
        
        # Set slide background color
        try:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = _hex_to_rgb(canvas_background)
        except Exception as bg_error:
            print(f"Failed to set background color: {bg_error}")
        
        # Add canvas border if specified
        canvas_border_width = layout_config.get('canvasBorderWidth', 0)
        canvas_border_color = layout_config.get('canvasBorderColor', '#000000')
        
        if canvas_border_width > 0:
            try:
                # Add border as a rectangle shape
                border_shape = slide.shapes.add_shape(
                    1,  # Rectangle auto shape
                    Inches(0), Inches(0),
                    _pixels_to_inches(layout_config.get('canvasWidth', 1123)),
                    _pixels_to_inches(layout_config.get('canvasHeight', 794))
                )
                
                # Configure border
                border_shape.fill.background()  # No fill, just border
                border_shape.line.color.rgb = _hex_to_rgb(canvas_border_color)
                border_shape.line.width = Pt(canvas_border_width)
                
                print(f"Added canvas border: {canvas_border_width}pt, color: {canvas_border_color}")
            except Exception as border_error:
                print(f"Failed to add canvas border: {border_error}")
        
        # Process each element from the canvas
        for i, element in enumerate(elements):
            try:
                print(f"Processing element {i+1}/{len(elements)}: {element.get('type', 'unknown')}")
                _add_element_to_slide(slide, element, unit, group_data)
            except Exception as element_error:
                print(f"Error processing element {i+1}: {element_error}")
                continue  # Skip problematic elements but continue with others
        
        print(f"Slide created successfully for unit: {unit.get('name')}")
        return slide
        
    except Exception as e:
        print(f"Error creating slide for unit {unit.get('name', 'Unknown')}: {e}")
        # Create a minimal slide with just the unit name as fallback
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add unit name as fallback
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
        text_frame = text_box.text_frame
        text_frame.text = f"{unit.get('name', 'Unknown Unit')} - {unit.get('unit_class', 'Unknown Class')}"
        
        return slide

def _create_grid_slide(prs: Presentation, units: List[Dict[str, Any]], group_data: Dict[str, Any], 
                      rows: int, cols: int, page_num: int) -> Any:
    """Create a grid slide with multiple units"""
    
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add page title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = f"{group_data.get('name', 'Gruppo')} - Pagina {page_num}"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(24)
    title_para.font.bold = True
    
    # Calculate grid dimensions
    slide_width = Inches(12.5)  # Leave margins
    slide_height = Inches(6)    # Leave space for title
    cell_width = slide_width / cols
    cell_height = slide_height / rows
    
    # Place units in grid
    for i, unit in enumerate(units):
        row = i // cols
        col = i % cols
        
        x = Inches(0.5) + col * cell_width
        y = Inches(1) + row * cell_height
        
        _add_unit_to_grid_cell(slide, unit, group_data, x, y, cell_width, cell_height)
    
    return slide

def _add_unit_to_grid_cell(slide: Any, unit: Dict[str, Any], group_data: Dict[str, Any], 
                          x: Inches, y: Inches, width: Inches, height: Inches):
    """Add a unit summary to a grid cell"""
    
    # Add cell border
    cell_box = slide.shapes.add_textbox(x, y, width, height)
    
    # Add unit name
    name_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.1), 
                                       width - Inches(0.2), Inches(0.4))
    name_frame = name_box.text_frame
    name_frame.text = unit.get('name', 'Nome Unità')
    name_para = name_frame.paragraphs[0]
    name_para.alignment = PP_ALIGN.CENTER
    name_para.font.size = Pt(14)
    name_para.font.bold = True
    
    # Add unit class
    class_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.5), 
                                        width - Inches(0.2), Inches(0.3))
    class_frame = class_box.text_frame
    class_frame.text = f"Classe: {unit.get('unit_class', 'N/A')}"
    class_para = class_frame.paragraphs[0]
    class_para.alignment = PP_ALIGN.CENTER
    class_para.font.size = Pt(10)
    
    # Try to add silhouette if available
    layout_config = unit.get('layout_config', {})
    elements = layout_config.get('elements', [])
    silhouette_element = next((el for el in elements if el.get('type') == 'silhouette' and el.get('image')), None)
    
    if silhouette_element and silhouette_element.get('image'):
        try:
            # Add silhouette image (simplified - would need proper image handling)
            img_y = y + Inches(0.8)
            img_height = height - Inches(1.0)
            # Note: Real implementation would need to handle image loading and sizing
        except Exception:
            pass

def _add_element_to_slide(slide: Any, element: Dict[str, Any], unit: Dict[str, Any], group_data: Dict[str, Any]):
    """Add a canvas element to the slide"""
    
    element_type = element.get('type')
    x = _pixels_to_inches(element.get('x', 0))
    y = _pixels_to_inches(element.get('y', 0))
    width = _pixels_to_inches(element.get('width', 100))
    height = _pixels_to_inches(element.get('height', 30))
    
    if element_type in ['text', 'unit_name', 'unit_class']:
        # Add text element
        text_box = slide.shapes.add_textbox(x, y, width, height)
        text_frame = text_box.text_frame
        text_frame.text = element.get('content', '')
        
        # Apply styling
        paragraph = text_frame.paragraphs[0]
        font = paragraph.font
        
        style = element.get('style', {})
        font.size = Pt(style.get('fontSize', 16))
        font.bold = style.get('fontWeight') == 'bold'
        
        # Text color
        color = style.get('color', '#000000')
        font.color.rgb = _hex_to_rgb(color)
        
        # Text alignment
        text_align = style.get('textAlign', 'left')
        if text_align == 'center':
            paragraph.alignment = PP_ALIGN.CENTER
        elif text_align == 'right':
            paragraph.alignment = PP_ALIGN.RIGHT
        
        # Apply background color and borders
        if style.get('backgroundColor'):
            text_box.fill.solid()
            text_box.fill.fore_color.rgb = _hex_to_rgb(style.get('backgroundColor'))
        
        # Apply borders
        border_width = style.get('borderWidth', 0)
        if border_width > 0:
            text_box.line.color.rgb = _hex_to_rgb(style.get('borderColor', '#000000'))
            text_box.line.width = Pt(border_width)
            border_style = style.get('borderStyle', 'solid')
            # PowerPoint supports different line styles, but solid is most compatible
        
        # Apply border radius (limited support in PowerPoint)
        border_radius = style.get('borderRadius', 0)
        if border_radius > 0:
            # PowerPoint has limited border radius support, this is approximate
            try:
                text_box.adjustments[0] = border_radius / 100.0  # Approximate conversion
            except:
                pass  # Ignore if adjustments not available
    
    elif element_type in ['logo', 'flag', 'silhouette']:
        # Handle images - with remote URL support
        image_path = element.get('image')
        if image_path:
            temp_files_to_cleanup = []
            try:
                # Apply group template overrides
                if element_type == 'logo' and group_data.get('override_logo') and group_data.get('template_logo_path'):
                    image_path = group_data.get('template_logo_path')
                elif element_type == 'flag' and group_data.get('override_flag') and group_data.get('template_flag_path'):
                    image_path = group_data.get('template_flag_path')
                
                print(f"Processing {element_type} image: {image_path}")
                
                # Get the actual image path (download if remote)
                actual_image_path = _get_image_path(image_path)
                
                if actual_image_path:
                    # If it was downloaded, mark for cleanup
                    if image_path.startswith('http'):
                        temp_files_to_cleanup.append(actual_image_path)
                    
                    try:
                        print(f"Adding {element_type} image to slide: {actual_image_path}")
                        print(f"Target dimensions: {width} x {height} at position ({x}, {y})")
                        
                        # Get image dimensions first using PIL to calculate proper aspect ratio
                        try:
                            with PILImage.open(actual_image_path) as img:
                                img_width, img_height = img.size
                                print(f"Original image size: {img_width} x {img_height}")
                                
                                # Calculate aspect ratios
                                img_aspect = img_width / img_height
                                target_aspect = width / height
                                
                                print(f"Image aspect: {img_aspect:.3f}, Target aspect: {target_aspect:.3f}")
                                
                                # Calculate final dimensions maintaining aspect ratio
                                if img_aspect > target_aspect:
                                    # Image is wider than target - fit to width
                                    final_width = width
                                    final_height = width / img_aspect
                                else:
                                    # Image is taller than target - fit to height
                                    final_height = height
                                    final_width = height * img_aspect
                                
                                # Center the image
                                final_x = x + (width - final_width) / 2
                                final_y = y + (height - final_height) / 2
                                
                                print(f"Final dimensions: {final_width} x {final_height} at ({final_x}, {final_y})")
                                
                                # Add image with calculated dimensions
                                picture = slide.shapes.add_picture(actual_image_path, final_x, final_y, final_width, final_height)
                                print(f"Successfully added {element_type} image with proper aspect ratio")
                                
                        except Exception as pil_error:
                            print(f"PIL processing failed: {pil_error}, using fallback method")
                            import traceback
                            traceback.print_exc()
                            # Fallback: add image at original size and position
                            try:
                                picture = slide.shapes.add_picture(actual_image_path, x, y, width, height)
                                print(f"Added {element_type} image with fallback method")
                            except Exception as fallback_error:
                                print(f"Fallback also failed: {fallback_error}")
                                # Last resort: create placeholder text
                                text_box = slide.shapes.add_textbox(x, y, width, height)
                                text_frame = text_box.text_frame
                                text_frame.text = f"[{element_type.upper()}]"
                                print(f"Created text placeholder for {element_type}")
                        
                        # Clean up temporary files
                        for temp_file in temp_files_to_cleanup:
                            try:
                                os.unlink(temp_file)
                                print(f"Cleaned up temporary file: {temp_file}")
                            except:
                                pass
                        
                    except Exception as img_error:
                        print(f"Error adding image {actual_image_path}: {img_error}")
                        # Clean up temporary files on error
                        for temp_file in temp_files_to_cleanup:
                            try:
                                os.unlink(temp_file)
                            except:
                                pass
                        # Add placeholder text if image fails
                        text_box = slide.shapes.add_textbox(x, y, width, height)
                        text_frame = text_box.text_frame
                        text_frame.text = f"[{element_type.upper()}]"
                else:
                    print(f"Could not get image for {element_type}: {image_path}")
                    
                    # For silhouettes, try to find in unit data
                    if element_type == 'silhouette' and unit:
                        # Try unit.silhouette_path
                        unit_silhouette = unit.get('silhouette_path')
                        if unit_silhouette:
                            print(f"🔍 Trying unit silhouette_path: {unit_silhouette}")
                            actual_image_path_alt = _get_image_path(unit_silhouette)
                            if actual_image_path_alt:
                                try:
                                    picture = slide.shapes.add_picture(actual_image_path_alt, x, y, width, height)
                                    print(f"✅ Added silhouette from unit.silhouette_path")
                                    return  # Success, exit function
                                except Exception as alt_error:
                                    print(f"❌ Failed to add unit silhouette: {alt_error}")
                    
                    # Add placeholder text if all methods fail
                    text_box = slide.shapes.add_textbox(x, y, width, height)
                    text_frame = text_box.text_frame
                    text_frame.text = f"[{element_type.upper()}]"
                    
                    # Style the placeholder better
                    text_frame.margin_left = Inches(0.1)
                    text_frame.margin_right = Inches(0.1)
                    text_frame.margin_top = Inches(0.1)
                    text_frame.margin_bottom = Inches(0.1)
                    
                    paragraph = text_frame.paragraphs[0]
                    paragraph.alignment = PP_ALIGN.CENTER
                    font = paragraph.font
                    font.size = Pt(12)
                    font.bold = True
                    font.color.rgb = RGBColor(128, 128, 128)  # Gray color
                    
                    # Add border to placeholder
                    text_box.line.color.rgb = RGBColor(200, 200, 200)
                    text_box.line.width = Pt(1)
                    text_box.fill.solid()
                    text_box.fill.fore_color.rgb = RGBColor(245, 245, 245)  # Light gray background
                    
                    print(f"Created styled placeholder for {element_type} - image file missing: {image_path}")
            except Exception as e:
                print(f"Exception processing {element_type} image: {e}")
                # Clean up temporary files on error
                for temp_file in temp_files_to_cleanup:
                    try:
                        os.unlink(temp_file)
                    except:
                        pass
                # If image fails, add placeholder text
                text_box = slide.shapes.add_textbox(x, y, width, height)
                text_frame = text_box.text_frame
                text_frame.text = f"[{element_type.upper()}]"
    
    elif element_type == 'table':
        # Add table data as a real PowerPoint table
        table_data = element.get('tableData', [])
        if table_data and len(table_data) > 0:
            try:
                rows = len(table_data)
                cols = len(table_data[0]) if rows > 0 else 2
                
                # Create PowerPoint table
                table_shape = slide.shapes.add_table(rows, cols, x, y, width, height)
                table = table_shape.table
                
                # Set table data
                for row_idx, row_data in enumerate(table_data):
                    for col_idx, cell_data in enumerate(row_data):
                        if col_idx < cols:  # Safety check
                            cell = table.cell(row_idx, col_idx)
                            cell.text = str(cell_data)
                            
                            # Style the cell
                            cell_font = cell.text_frame.paragraphs[0].font
                            cell_font.size = Pt(9)
                            
                            # Style header row differently
                            if row_idx == 0:
                                cell_font.bold = True
                                cell_font.color.rgb = RGBColor(255, 255, 255)  # White text
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = RGBColor(31, 73, 125)  # Dark blue background
                            else:
                                # Alternate row colors for better readability
                                if row_idx % 2 == 0:
                                    cell.fill.solid()
                                    cell.fill.fore_color.rgb = RGBColor(242, 242, 242)  # Light gray
                                else:
                                    cell.fill.solid()
                                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
                
                # Set table style
                table_shape.table.first_row = True  # Enable header row formatting
                
                print(f"Created PowerPoint table with {rows} rows and {cols} columns")
                
            except Exception as table_error:
                print(f"Error creating PowerPoint table: {table_error}")
                # Fallback to text if table creation fails
                text_content = []
                for row in table_data:
                    text_content.append(' | '.join(str(cell) for cell in row))
                
                text_box = slide.shapes.add_textbox(x, y, width, height)
                text_frame = text_box.text_frame
                text_frame.text = '\n'.join(text_content)
                
                # Style as table-like
                paragraph = text_frame.paragraphs[0]
                font = paragraph.font
                font.size = Pt(10)

def _pixels_to_inches(pixels: float) -> Inches:
    """Convert pixels to inches using current scale"""
    global CANVAS_TO_PPT_SCALE
    scaled_pixels = pixels * CANVAS_TO_PPT_SCALE
    return Inches(scaled_pixels / 120.0)  # Updated to match new DPI

def _download_image(image_url: str) -> Optional[str]:
    """Download image from URL and return temporary file path"""
    try:
        print(f"Downloading image from: {image_url}")
        
        # Set headers to mimic a browser request
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        
        response = requests.get(image_url, headers=headers, timeout=10)
        response.raise_for_status()
        
        # Get file extension from URL or default to .png
        file_ext = '.png'
        if '.' in image_url.split('/')[-1]:
            file_ext = '.' + image_url.split('.')[-1].split('?')[0]  # Remove query params
        
        # Create temporary file in server directory
        temp_dir = "./data/temp"
        os.makedirs(temp_dir, exist_ok=True)
        
        # Generate unique filename
        import uuid
        temp_filename = f"download_{uuid.uuid4().hex}{file_ext}"
        temp_path = os.path.join(temp_dir, temp_filename)
        
        # Write downloaded content
        with open(temp_path, 'wb') as f:
            f.write(response.content)
        
        print(f"Image downloaded successfully to: {temp_path}")
        return temp_path
        
    except Exception as e:
        print(f"Error downloading image {image_url}: {e}")
        return None

def _get_image_path(image_path: str) -> Optional[str]:
    """Get the correct image path, downloading remote images if needed"""
    print(f"🔍 _get_image_path called with: {image_path}")
    
    if not image_path:
        print("❌ Empty image_path provided")
        return None
    
    # If it's base64 data, convert to temporary file
    if image_path.startswith('data:image/'):
        print(f"🔄 Converting base64 image to temporary file")
        return _convert_base64_to_temp_file(image_path)
    
    # If it's an HTTP URL, download it
    if image_path.startswith('http://') or image_path.startswith('https://'):
        print(f"📥 Downloading remote image: {image_path}")
        return _download_image(image_path)
    
    # Handle local file paths
    if image_path.startswith('/api/static/'):
        local_path = image_path.replace('/api/static/', './data/uploads/')
        print(f"🔄 Converted /api/static/ path: {image_path} -> {local_path}")
    elif image_path.startswith('../data/uploads/'):
        local_path = image_path.replace('../data/uploads/', './data/uploads/')
        print(f"🔄 Converted ../ path: {image_path} -> {local_path}")
    elif not image_path.startswith('./data/uploads/'):
        local_path = f"./data/uploads/{image_path}"
        print(f"🔄 Added data/uploads prefix: {image_path} -> {local_path}")
    else:
        local_path = image_path
        print(f"✅ Using path as-is: {local_path}")
    
    # Check if local file exists
    print(f"🔍 Checking if file exists: {local_path}")
    if os.path.exists(local_path):
        print(f"✅ File found: {local_path}")
        return local_path
    else:
        print(f"❌ Local image file not found: {local_path}")
        
        # Try different path variations
        variations = [
            image_path,
            os.path.join('./data/uploads', os.path.basename(image_path)),
            os.path.join('./data/uploads/silhouettes', os.path.basename(image_path)),
            os.path.join('./data/uploads/logos', os.path.basename(image_path)),
            os.path.join('./data/uploads/flags', os.path.basename(image_path))
        ]
        
        for variation in variations:
            print(f"🔍 Trying variation: {variation}")
            if os.path.exists(variation):
                print(f"✅ Found at variation: {variation}")
                return variation
        
        print(f"❌ No variations found for: {image_path}")
        return None

def _convert_base64_to_temp_file(base64_data: str) -> Optional[str]:
    """Convert base64 image data to temporary file"""
    try:
        print(f"🔄 Converting base64 to temp file")
        
        # Extract the base64 data
        if ',' in base64_data:
            header, data = base64_data.split(',', 1)
            # Extract format from header (e.g., data:image/png;base64)
            if 'image/' in header:
                image_format = header.split('image/')[1].split(';')[0]
            else:
                image_format = 'png'  # default
        else:
            data = base64_data
            image_format = 'png'
        
        # Decode base64
        image_data = base64.b64decode(data)
        
        # Create temporary file in server directory
        temp_dir = "./data/temp"
        os.makedirs(temp_dir, exist_ok=True)
        
        # Generate unique filename
        import uuid
        temp_filename = f"img_{uuid.uuid4().hex}.{image_format}"
        temp_path = os.path.join(temp_dir, temp_filename)
        
        # Write image data
        with open(temp_path, 'wb') as f:
            f.write(image_data)
        
        print(f"✅ Base64 converted to temp file: {temp_path}")
        return temp_path
        
    except Exception as e:
        print(f"❌ Failed to convert base64 to temp file: {e}")
        return None

def _hex_to_rgb(hex_color: str):
    """Convert hex color to RGB"""
    hex_color = hex_color.lstrip('#')
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )